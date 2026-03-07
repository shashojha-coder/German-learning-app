#!/usr/bin/env python3
"""
Restyle Q4Q5 to match MDA Rollout format, and add APA citations
to MDA Rollout (references slide + speaker notes on every slide).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from copy import deepcopy
import copy

# ─── APA REFERENCES ──────────────────────────────────────────────────────────

APA_REFERENCES = [
    '[1]  PMI. (2021). PMBOK® Guide — A guide to the project management body of knowledge (7th ed.). Project Management Institute.',
    '[2]  Kerzner, H. (2022). Project management: A systems approach to planning, scheduling, and controlling (13th ed.). Wiley.',
    '[4]  Galbraith, J. R. (1971). Matrix organization designs: How to combine functional and project forms. Business Horizons, 14(1), 29–40.',
    '[5]  Davis, S. M., & Lawrence, P. R. (1977). Matrix. Addison-Wesley.',
    '[7]  Turner, J. R., & Cochrane, R. A. (1993). Goals-and-methods matrix: Coping with projects with ill-defined goals and/or methods of achieving them. International Journal of Project Management, 11(2), 93–102.',
    '[8]  Andersen, E. S., Grude, K. V., & Haug, T. (2009). Goal directed project management (4th ed.). Kogan Page. Cited in Miranda, E. (2019). Milestone planning. PM World Journal, VIII(XI).',
    '[9]  Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54.',
    '[10] Cooper, R. G. (2008). The Stage-Gate® idea-to-launch process — Update, what\'s new, and NexGen systems. Journal of Product Innovation Management, 25(3), 213–232.',
]

# Speaker notes per slide for MDA Rollout (slide index 0–8)
ROLLOUT_SPEAKER_NOTES = {
    0: (  # Slide 1 — Title
        "This presentation covers the MDA strategic timeline and rollout plan "
        "for the new DB FZI Cottbus plant, targeting a June 2027 Start of Production (SoP). "
        "The plan follows a Stage-Gate® methodology and tiered system prioritisation.\n\n"
        "References:\n"
        "PMI. (2021). PMBOK® Guide — A guide to the project management body of knowledge (7th ed.). Project Management Institute.\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54."
    ),
    1: (  # Slide 2 — Agenda
        "The agenda covers five key areas: hybrid rollout strategy, system prioritisation tiers, "
        "the four-phase timeline to SoP, key tools & methods, and the MDA Blueprint for cross-plant rollout.\n\n"
        "References:\n"
        "Kerzner, H. (2022). Project management: A systems approach to planning, scheduling, and controlling (13th ed.). Wiley.\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54."
    ),
    2: (  # Slide 3 — Rollout Strategy (Hybrid Sequencing)
        "The hybrid sequencing logic combines three criteria: system criticality (safety-critical first), "
        "building section readiness (completed areas first), and interface type clustering (OPC-UA bundles). "
        "This avoids the limitations of a purely system-type or section-only approach.\n\n"
        "References:\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54, p. 47.\n"
        "Kerzner, H. (2022). Project management: A systems approach to planning, scheduling, and controlling (13th ed.). Wiley, pp. 108–133.\n"
        "Turner, J. R., & Cochrane, R. A. (1993). Goals-and-methods matrix: Coping with projects with ill-defined goals and/or methods of achieving them. International Journal of Project Management, 11(2), 93–102."
    ),
    3: (  # Slide 4 — System Prioritisation (Three Tiers)
        "The 100 systems are categorised into three tiers based on operational criticality. "
        "Tier 1 (~15–20 systems) covers safety-critical assets such as wheel lathes, lifting systems, "
        "and brake test rigs. Tier 2 (~50 systems) addresses OEE-impact assets. Tier 3 (~30–35 systems) "
        "covers standard workshop equipment. This prioritisation follows Cooper's (1990) "
        "recommendation to sequence deliverables by value and risk.\n\n"
        "References:\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54.\n"
        "PMI. (2021). PMBOK® Guide (7th ed.). Project Management Institute, p. 245 — definition of milestones.\n"
        "Turner, J. R., & Cochrane, R. A. (1993). Goals-and-methods matrix. International Journal of Project Management, 11(2), 93–102."
    ),
    4: (  # Slide 5 — Rollout Gantt Chart
        "This Gantt chart shows the tier-by-tier rollout schedule from Q1 2026 to the June 2027 SoP. "
        "Phase 0 (Foundation) establishes architecture and vendor contracts. Phase 1 (Pilot) connects "
        "Tier 1 safety-critical systems. Phase 2 (Scale-Up) rolls out Tier 2 and Tier 3 systems. "
        "Phase 3 (SoP Ready) completes acceptance and handover. A 4–6 week buffer is embedded "
        "at each phase transition to absorb construction or vendor delays.\n\n"
        "References:\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54, p. 47 — phase-gate scheduling.\n"
        "Kerzner, H. (2022). Project management (13th ed.). Wiley, p. 133 — CPM scheduling standard for hard deadlines.\n"
        "Cooper, R. G. (2008). The Stage-Gate® idea-to-launch process — Update, what's new, and NexGen systems. Journal of Product Innovation Management, 25(3), 213–232."
    ),
    5: (  # Slide 6 — Four-Phase Timeline
        "The four phases follow Cooper's (1990, 2008) Stage-Gate® model: each phase ends with a "
        "management decision gate before the next phase is authorised. Phase 0 finalises the MDA "
        "architecture, Phase 1 connects pilot Tier 1 systems, Phase 2 scales to remaining ~80 systems, "
        "and Phase 3 delivers full acceptance and SoP handover.\n\n"
        "References:\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54.\n"
        "Cooper, R. G. (2008). The Stage-Gate® idea-to-launch process. Journal of Product Innovation Management, 25(3), 213–232, p. 215.\n"
        "PMI. (2021). PMBOK® Guide (7th ed.). Project Management Institute — milestone-based project planning.\n"
        "Andersen, E. S., Grude, K. V., & Haug, T. (2009). Goal directed project management (4th ed.). Kogan Page — participatory milestone planning."
    ),
    6: (  # Slide 7 — Key Tools & Methods
        "The key tools split into project management (MS Project/Jira, RACI Matrix, Risk Register, "
        "Interface Catalogue) and technical infrastructure (IoT Middleware, OPC-UA Adapters, "
        "OEE Dashboard, MES/SAP-PM Connectors). MS Project with CPM is the scheduling standard "
        "for projects with hard external deadlines (Kerzner, 2022, p. 133).\n\n"
        "References:\n"
        "Kerzner, H. (2022). Project management: A systems approach to planning, scheduling, and controlling (13th ed.). Wiley, p. 133.\n"
        "PMI. (2021). PMBOK® Guide (7th ed.). Project Management Institute, p. 160 — RACI framework.\n"
        "Davis, S. M., & Lawrence, P. R. (1977). Matrix. Addison-Wesley — matrix governance and role clarity."
    ),
    7: (  # Slide 8 — Cross-Plant Strategy (Blueprint)
        "The MDA Blueprint captures all technical and organisational learnings from Cottbus in a "
        "structured, transferable format. Blueprint v1 is delivered at the end of Phase 1 (Q4 2026), "
        "and Blueprint v2 — the definitive cross-plant guide — at SoP (June 2027). Components include "
        "interface connector specs, data model standards, vendor templates, and a Works Council playbook.\n\n"
        "References:\n"
        "Cooper, R. G. (2008). The Stage-Gate® idea-to-launch process. Journal of Product Innovation Management, 25(3), 213–232 — knowledge capture across gates.\n"
        "Kerzner, H. (2022). Project management (13th ed.). Wiley — lessons learned and knowledge management.\n"
        "PMI. (2021). PMBOK® Guide (7th ed.). Project Management Institute — organisational process assets."
    ),
    8: (  # Slide 9 — Key Takeaways
        "Four key takeaways: (1) Hybrid sequencing combines criticality, readiness, and interface "
        "clustering. (2) Four buffer-protected phases lead to the June 2027 SoP. (3) Tiered "
        "prioritisation ensures Tier 1 safety-critical systems deliver value first. (4) The MDA "
        "Blueprint transforms Cottbus learnings into reusable assets for all future DB FZI plants.\n\n"
        "References:\n"
        "PMI. (2021). PMBOK® Guide — A guide to the project management body of knowledge (7th ed.). Project Management Institute.\n"
        "Kerzner, H. (2022). Project management: A systems approach to planning, scheduling, and controlling (13th ed.). Wiley.\n"
        "Cooper, R. G. (1990). Stage-gate systems: A new tool for managing new products. Business Horizons, 33(3), 44–54.\n"
        "Cooper, R. G. (2008). The Stage-Gate® idea-to-launch process — Update, what's new, and NexGen systems. Journal of Product Innovation Management, 25(3), 213–232.\n"
        "Turner, J. R., & Cochrane, R. A. (1993). Goals-and-methods matrix. International Journal of Project Management, 11(2), 93–102.\n"
        "Galbraith, J. R. (1971). Matrix organization designs. Business Horizons, 14(1), 29–40.\n"
        "Davis, S. M., & Lawrence, P. R. (1977). Matrix. Addison-Wesley.\n"
        "Andersen, E. S., Grude, K. V., & Haug, T. (2009). Goal directed project management (4th ed.). Kogan Page."
    ),
}


# ═══════════════════════════════════════════════════════════════════════════════
#  PART 1 — Restyle Q4Q5 to match Rollout slide dimensions & backgrounds
# ═══════════════════════════════════════════════════════════════════════════════

def restyle_q4q5():
    """Scale Q4Q5 from 14630400×8229600 to 9144000×5143500 and fix backgrounds."""
    src = "Q4Q5_Final_9slides.pptx"
    dst = "Q4Q5_Final_9slides_restyled.pptx"

    prs = Presentation(src)

    target_w = 9144000
    target_h = 5143500
    scale = target_w / prs.slide_width  # 0.625

    prs.slide_width = target_w
    prs.slide_height = target_h

    for i, slide in enumerate(prs.slides):
        # Fix backgrounds — make all slides cream like rollout
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xF0)

        # Scale all shapes
        for shape in slide.shapes:
            shape.left = int(shape.left * scale)
            shape.top = int(shape.top * scale)
            shape.width = int(shape.width * scale)
            shape.height = int(shape.height * scale)

            # Scale font sizes within text frames
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size is not None:
                            run.font.size = int(run.font.size * scale)

    prs.save(dst)
    print(f"  ✓ Saved restyled Q4Q5 → {dst}")
    print(f"    New dimensions: {target_w} x {target_h} (matched to Rollout)")
    print(f"    All slides: background → #FFF8F0 (cream)")
    print(f"    Scale factor: {scale:.4f}")


# ═══════════════════════════════════════════════════════════════════════════════
#  PART 2 — Add APA reference slide + speaker notes to MDA Rollout
# ═══════════════════════════════════════════════════════════════════════════════

def add_citations_to_rollout():
    """Add a references slide at the end and APA speaker notes to every slide."""
    src = "MDA_Rollout_Restyled (1).pptx"
    dst = "MDA_Rollout_Restyled_Citations.pptx"

    prs = Presentation(src)
    slide_w = prs.slide_width   # 9144000
    slide_h = prs.slide_height  # 5143500

    # --- Add speaker notes to each existing slide ---
    for idx, slide in enumerate(prs.slides):
        if idx in ROLLOUT_SPEAKER_NOTES:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            tf.clear()
            tf.text = ROLLOUT_SPEAKER_NOTES[idx]

    # --- Add APA References slide at the end ---
    # Use the blank layout (last layout, or find one without content placeholders)
    layout = prs.slide_layouts[len(prs.slide_layouts) - 1]  # typically blank
    ref_slide = prs.slides.add_slide(layout)

    # Set cream background
    ref_slide.background.fill.solid()
    ref_slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xF0)

    # Header banner (dark teal bar matching rollout header style)
    from pptx.util import Emu
    banner = ref_slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Emu(228600), Emu(228600),
        Emu(slide_w - 457200), Emu(502920)
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0x2C, 0x3F, 0x42)
    banner.line.fill.background()

    # Banner subtitle
    from pptx.util import Pt
    sub_box = ref_slide.shapes.add_textbox(
        Emu(342900), Emu(246888),
        Emu(slide_w - 685800), Emu(160020)
    )
    sub_tf = sub_box.text_frame
    sub_p = sub_tf.paragraphs[0]
    sub_run = sub_p.add_run()
    sub_run.text = "APA REFERENCES"
    sub_run.font.name = "Open Sans"
    sub_run.font.size = Pt(8.5)
    sub_run.font.bold = True
    sub_run.font.color.rgb = RGBColor(0xE2, 0xC8, 0xB5)

    # Banner title
    title_box = ref_slide.shapes.add_textbox(
        Emu(342900), Emu(388620),
        Emu(slide_w - 685800), Emu(320040)
    )
    title_tf = title_box.text_frame
    title_p = title_tf.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = "Full Reference List"
    title_run.font.name = "Bitter"
    title_run.font.size = Pt(24)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # References body
    ref_top = 880000
    ref_spacing = 340000

    for i, ref_text in enumerate(APA_REFERENCES):
        y = ref_top + i * ref_spacing
        tb = ref_slide.shapes.add_textbox(
            Emu(342900), Emu(y),
            Emu(slide_w - 685800), Emu(ref_spacing)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = ref_text
        run.font.name = "Roboto"
        run.font.size = Pt(9.5)
        run.font.color.rgb = RGBColor(0x2B, 0x2E, 0x3C)

    # Footer bar
    footer = ref_slide.shapes.add_shape(
        1,  # RECTANGLE
        Emu(0), Emu(slide_h - 263525),
        Emu(slide_w), Emu(263525)
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(0x2C, 0x3F, 0x42)
    footer.line.fill.background()

    footer_tb = ref_slide.shapes.add_textbox(
        Emu(0), Emu(slide_h - 263525),
        Emu(slide_w), Emu(263525)
    )
    ft_tf = footer_tb.text_frame
    ft_p = ft_tf.paragraphs[0]
    ft_p.alignment = PP_ALIGN.CENTER
    ft_run = ft_p.add_run()
    ft_run.text = "DB FZI — MDA Project  ·  New Cottbus Plant  ·  Target: June 2027 SoP"
    ft_run.font.name = "Open Sans"
    ft_run.font.size = Pt(10)
    ft_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # Add speaker notes to references slide
    if not ref_slide.has_notes_slide:
        ref_slide.notes_slide  # accessing it creates it
    ref_notes = ref_slide.notes_slide
    if ref_notes and ref_notes.notes_text_frame:
        ref_notes.notes_text_frame.text = (
            "This slide contains the full APA-formatted reference list for all academic "
            "and professional sources cited throughout the MDA Rollout presentation.\n\n"
            + "\n".join(APA_REFERENCES)
        )

    prs.save(dst)
    print(f"  ✓ Saved Rollout with citations → {dst}")
    print(f"    Added speaker notes to all {len(prs.slides) - 1} original slides")
    print(f"    Added References slide (slide {len(prs.slides)})")
    print(f"    {len(APA_REFERENCES)} APA citations included")


# ═══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 60)
    print("  PART 1: Restyling Q4Q5 to match Rollout format")
    print("=" * 60)
    restyle_q4q5()

    print()
    print("=" * 60)
    print("  PART 2: Adding APA citations to MDA Rollout")
    print("=" * 60)
    add_citations_to_rollout()

    print()
    print("Done! Two new files created:")
    print("  1. Q4Q5_Final_9slides_restyled.pptx")
    print("  2. MDA_Rollout_Restyled_Citations.pptx")
